import pyautogui as pg
import cv2
import numpy as np
import tkinter as tk
from paddleocr import PaddleOCR,draw_ocr
from PIL import Image
import keyboard

import  os
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

def ppt(size,boxes,txts):
    # 实例化 ppt 文档对象
    try:
        prs = Presentation('test.pptx')
    except:
        prs = Presentation()

    # 插入幻灯片
    blank_slide = prs.slide_layouts[6]
    slide_1 = prs.slides.add_slide(blank_slide)
    for i in range(len(boxes)):
    # 预设位置及大小

        left =Cm(boxes[i][0][0]/size[2]*25.4)  # left，top为相对位置
        top = Cm(boxes[i][0][1]/size[3]*19.05*0.9) 
        width = Cm((boxes[i][1][0]-boxes[i][0][0])/size[2]*25.4)   # width，height为文本框的大小
        height = Cm((boxes[i][3][1]-boxes[i][0][1])/size[3]*19.05)  

        # 在指定位置添加文本框
        textbox = slide_1.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        
        tf.margin_left=0
        tf.margin_right=0 
        tf.margin_top=0 
        tf.margin_bottom=0

        # 在文本框中写入文字

        para = tf.add_paragraph()    # 新增段落

        para.text = txts[i]  # 向段落写入文字
        para.alignment = PP_ALIGN.LEFT    # 居中
        para.line_spacing = 0.0    # 0.0 倍的行距

        ### 设置字体
        font = para.font
        font.name = '微软雅黑'    # 字体类型
        #font.name = 'Calibri'    # 字体类型       
        font.bold = False    # 加粗
        font.size = Pt(int((boxes[i][2][1]-boxes[i][0][1])*0.75))   # 大小
        

        # 保存 ppt
        prs.save('test.pptx')


ocr = PaddleOCR(use_angle_cls=True, lang="ch") 





class App:
    def __init__(self, root):
        root.title("盗幕笔记")
        root.bind("<Control-x>",self.ocr)
        frame = tk.Frame(root)
        
        frame.pack()
        self.size=[0,0,100,100]
        self.hi_there = tk.Button(frame, text="记笔记", fg="blue", command=self.ocr)
        
        self.hi_there.pack(side=tk.LEFT)    
        self.hi_there = tk.Button(frame, text="截图区域设置", fg="blue", command=self.pos)
        self.hi_there.pack(side=tk.RIGHT) 


    def pos(self):
        print('请用ctrl+1确定左上角位置：')
        while True:

            hotkey = keyboard.read_hotkey()
            if 'ctrl' in hotkey and '1' in hotkey:

                x1, y1 = pg.position()
                print('左上角位置：',[x1,y1]) 
                keyboard._pressed_events={}
                break
        print('\n请用ctrl+2确定右下角位置：')                
        while True: 
            
            hotkey2 = keyboard.read_hotkey()

            if '2' in hotkey2 and 'ctrl' in hotkey2:
                x2, y2 = pg.position()
                print('右下角位置：',[x2,y2]) 
                keyboard._pressed_events={}
                break
  
        self.size=[x1,y1,x2-x1,y2-y1]
        return self.size

        
    def ocr(self,*event):
        im=pg.screenshot(region=self.size)
        img=cv2.cvtColor(np.asarray(im),cv2.COLOR_BGR2RGB)
        cv2.imwrite('img.jpg',img)
        img_path = 'img.jpg'
        result = ocr.ocr(img_path, cls=True)
        f = "biji.txt"
        with open(f,"a") as file: 
            for line in result:
                file.write(line[1][0] +"\n")
                print(line[1][0])
            file.write("\n")
        image = Image.open(img_path).convert('RGB')
        boxes = [line[0] for line in result]
        txts = [line[1][0] for line in result]

        scores = [line[1][1] for line in result]
        im_show = draw_ocr(image, boxes, txts, scores, font_path='./fonts/simfang.ttf')
        im_show = Image.fromarray(im_show)
        im_show.save('result.jpg')
        ppt(self.size,boxes,txts)

root = tk.Tk()
root.geometry("200x100+200+50")
root.attributes("-topmost", True)

app = App(root)

root.mainloop()

